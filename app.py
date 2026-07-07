import streamlit as st
import pandas as pd
from pptx import Presentation
import io
import zipfile

# 1. PARAMÉTRAGE DE L'INTERFACE WEB POUR STREAMLIT
st.set_page_config(page_title="Générateur de Rapports PPTX", page_icon="🚀", layout="centered")

st.title("Générateur Automatique de Rapports PPTX")
st.write("Cette interface permet de générer des présentations PowerPoint personnalisées par lot à partir d'un modèle et d'un tableau de données.")

# Zone d'explications pour l'utilisateur à distance
st.info("""
**Comment tester ?**
1. Téléversez un modèle PowerPoint (.pptx) contenant des balises comme `{{NOM_CLIENT}}`, `{{CHIFFRE_AFFAIRES}}`, ou `{{OBJECTIF}}`.
2. Téléversez votre fichier de données au format CSV.
3. Cliquez sur 'Lancer la génération' pour récupérer vos rapports personnalisés !
""")

# 2. MODULES DE TÉLÉCHARGEMENT (UPLOAD)
uploaded_template = st.file_uploader("1. Choisissez votre fichier Template PPTX", type="pptx")
uploaded_csv = st.file_uploader("2. Choisissez votre fichier de données CSV", type="csv")

# 3. MOTEUR D'INJECTION ET GÉNÉRATION
if uploaded_template and uploaded_csv:
    if st.button("Lancer la génération des rapports"):
        
        # Lecture du tableau CSV grâce à Pandas
        df = pd.read_csv(uploaded_csv)
        
        # Création d'un dossier ZIP virtuel en mémoire pour stocker tous les PPTX
        zip_buffer = io.BytesIO()
        
        with zipfile.ZipFile(zip_buffer, "w") as zip_file:
            
            # La Boucle : On traite le tableau ligne par ligne
            for index, row in df.iterrows():
    
                # On recharge une copie propre du template depuis la mémoire
                template_bytes = uploaded_template.getvalue()
                prs = Presentation(io.BytesIO(template_bytes))
    
            # --- CALCULS EN PYTHON ---
            ca = float(row['CA'])
            objectif = float(row['Objectif'])
    
            # 1. Calcul du taux de réussite
            taux_calculé = (ca / objectif) * 100
            taux_str = f"{taux_calculé:.1f} %"
    
            # 2. Logique pour la recommandation stratégique
            if taux_calculé >= 100:
                reco_str = "Objectifs atteints. Recommandation : Stratégie de fidélisation, développement de nouveaux comptes et phase d'up-selling."
                analyse_ia = "Excellente performance sur la période. La dynamique commerciale est solide et dépasse les prévisions."
            else:
                reco_str = "Objectifs non atteints. Recommandation : Mise en place d'un plan de relance commercial, analyse des freins à la conversion et accompagnement ciblé."
                analyse_ia = "Les résultats sont en deçà des attentes. Un ajustement stratégique est requis pour redresser la barre au prochain trimestre."

    # Parcours de toutes les slides et de toutes les formes textuelles
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        
                        # REMPLACEMENT DYNAMIQUE DES ANCIENNES ET NOUVELLES BALISES
                        if "{{NOM_CLIENT}}" in run.text:
                            run.text = run.text.replace("{{NOM_CLIENT}}", str(row['Client']))
                        if "{{CHIFFRE_AFFAIRES}}" in run.text:
                            run.text = run.text.replace("{{CHIFFRE_AFFAIRES}}", f"{ca:,.0f} €")
                        if "{{OBJECTIF}}" in run.text:
                            run.text = run.text.replace("{{OBJECTIF}}", f"{objectif:,.0f} €")
                        if "{{TAUX_REUSSITE}}" in run.text:
                            run.text = run.text.replace("{{TAUX_REUSSITE}}", taux_str)
                        if "{{ANALYSE_IA}}" in run.text:
                            run.text = run.text.replace("{{ANALYSE_IA}}", analyse_ia)
                        if "{{RECOMMANDATION_STRATEGIQUE}}" in run.text:
                            run.text = run.text.replace("{{RECOMMANDATION_STRATEGIQUE}}", reco_str)
        
        # 4. DISPOSITIF DE TÉLÉCHARGEMENT DU RÉSULTAT
        st.success("Tous les rapports ont été générés avec succès !")
        
        st.download_button(
            label="⬇Télécharger le package de rapports (ZIP)",
            data=zip_buffer.getvalue(),
            file_name="tous_les_rapports_generes.zip",
            mime="application/zip"
        )
